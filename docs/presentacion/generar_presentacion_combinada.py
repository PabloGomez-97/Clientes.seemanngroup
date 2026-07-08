"""
Genera una presentación ejecutiva RESUMIDA (10 diapositivas) que abarca ambos
portales de Seemann Group: Portal de Clientes y Portales Internos.
Enfoque general y de alto nivel para gerencia.

Uso: python generar_presentacion_combinada.py
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).parent
ASSETS_DIR = BASE_DIR / "assets"
OUTPUT = BASE_DIR / "portal-seemann-resumen-gerencia.pptx"
LOGO = ASSETS_DIR / "logocompleto.png"

PRIMARY = RGBColor(0xFF, 0x62, 0x00)
SECONDARY = RGBColor(0x23, 0x2F, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x23, 0x2F, 0x3E)
MUTED = RGBColor(0x66, 0x6E, 0x7A)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF8)
SLIDE_W = Inches(13.333)


def _set_run(run, *, size=18, bold=False, color=DARK_TEXT, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def _title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set_run(r, size=32, bold=True, color=SECONDARY)
    if subtitle:
        r2 = tf.add_paragraph().add_run()
        r2.text = subtitle
        _set_run(r2, size=16, color=MUTED)


def _bullets(slide, items, left=0.65, top=1.7, width=12, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = "•  " + item
        _set_run(r, size=size, color=DARK_TEXT)


def _notes(slide, notes: str):
    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    _accent_bar(slide)
    return slide


def _card(slide, left, top, w, h, title, desc, *, title_size=18, desc_size=13):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = PRIMARY
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(10)
    r = tf.paragraphs[0].add_run()
    r.text = title
    _set_run(r, size=title_size, bold=True, color=PRIMARY)
    if desc:
        r2 = tf.add_paragraph().add_run()
        r2.text = desc
        _set_run(r2, size=desc_size, color=DARK_TEXT)
    return card


# ── Slide 1: Portada ─────────────────────────────────────────────────────────
def s_cover(prs):
    slide = _blank(prs)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.8), Inches(0.9), width=Inches(3.2))
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(2.5))
    tf = box.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "Plataforma Digital Seemann Group"
    _set_run(r, size=40, bold=True, color=SECONDARY)
    r2 = tf.add_paragraph().add_run()
    r2.text = "Un ecosistema que conecta a clientes, equipo interno y proveedores"
    _set_run(r2, size=22, color=PRIMARY)
    r3 = tf.add_paragraph().add_run()
    r3.text = date.today().strftime("%B %Y").capitalize()
    _set_run(r3, size=14, color=MUTED)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.6), SLIDE_W, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = SECONDARY
    band.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11), Inches(0.5))
    tr = tb.text_frame.paragraphs[0].add_run()
    tr.text = "Resumen ejecutivo para gerencia"
    _set_run(tr, size=14, color=WHITE)
    _notes(
        slide,
        "Presentar la plataforma como un único ecosistema digital con dos caras: el portal "
        "para clientes y los portales para el equipo interno y proveedores. Visión general en 10 minutos.",
    )


# ── Slide 2: El ecosistema de un vistazo ─────────────────────────────────────
def s_ecosystem(prs):
    slide = _blank(prs)
    _title(slide, "El ecosistema de un vistazo", "Tres actores, una sola información compartida")
    _card(slide, 0.65, 2.0, 3.9, 2.6, "Clientes",
          "Cotizan, rastrean y revisan su cuenta en línea, 24/7", title_size=20, desc_size=14)
    _card(slide, 4.75, 2.0, 3.9, 2.6, "Equipo Seemann",
          "Cotiza, gestiona clientes, tarifas y reportes", title_size=20, desc_size=14)
    _card(slide, 8.85, 2.0, 3.8, 2.6, "Proveedores",
          "Publican y actualizan sus tarifas para Seemann", title_size=20, desc_size=14)
    arrow = slide.shapes.add_textbox(Inches(0.65), Inches(5.0), Inches(12), Inches(0.8))
    r = arrow.text_frame.paragraphs[0].add_run()
    r.text = "La misma información fluye entre los tres, sin correos ni planillas sueltas."
    _set_run(r, size=16, bold=True, color=SECONDARY)
    _notes(
        slide,
        "Idea central: todo se apoya en la misma base de información. Lo que el proveedor carga, "
        "el equipo lo cura, y el cliente lo consume. Menos fricción y menos errores manuales.",
    )


# ── Slide 3: Los dos portales ────────────────────────────────────────────────
def s_two_portals(prs):
    slide = _blank(prs)
    _title(slide, "Dos portales, un propósito común")
    left = _card(slide, 0.65, 1.9, 5.9, 3.7, "Portal de Clientes", "", title_size=24)
    tf = left.text_frame
    for line in [
        "Autoservicio para importadores y exportadores",
        "Cotizar, rastrear envíos y ver documentos",
        "Reportes financieros y operacionales",
        "Asistente virtual y contacto con su ejecutivo",
    ]:
        p = tf.add_paragraph()
        rr = p.add_run()
        rr.text = "•  " + line
        _set_run(rr, size=14, color=DARK_TEXT)

    right = _card(slide, 6.75, 1.9, 5.9, 3.7, "Portales Internos", "", title_size=24)
    tf2 = right.text_frame
    for line in [
        "Para el equipo Seemann y los proveedores",
        "Cotizar y gestionar la cartera de clientes",
        "Administrar tarifas y proveedores",
        "Reportes de negocio y control (dirección)",
    ]:
        p = tf2.add_paragraph()
        rr = p.add_run()
        rr.text = "•  " + line
        _set_run(rr, size=14, color=DARK_TEXT)
    _notes(
        slide,
        "El portal de clientes es la cara pública; los internos son la trastienda operativa. "
        "Accesos totalmente separados y seguros.",
    )


# ── Slide 4: Qué hace el cliente ─────────────────────────────────────────────
def s_client(prs):
    slide = _blank(prs)
    _title(slide, "Portal de Clientes", "Lo que el cliente puede hacer solo, sin llamar")
    items = [
        ("Cotizar", "Aéreo, marítimo y terrestre al instante"),
        ("Rastrear", "Seguimiento de envíos en tiempo real"),
        ("Documentos", "Descarga directa de sus archivos"),
        ("Reportes", "Estado financiero y operacional"),
        ("Alertas", "Avisos por cambios o retrasos"),
        ("Asistente IA", "Respuestas 24/7 dentro del portal"),
    ]
    for i, (t, d) in enumerate(items):
        row, col = divmod(i, 3)
        _card(slide, 0.65 + col * 4.15, 1.9 + row * 1.75, 3.8, 1.5, t, d)
    _notes(
        slide,
        "Resumir sin entrar en detalle: el cliente gana autonomía y visibilidad. "
        "Cada tarjeta es un servicio que antes requería un correo o llamada al ejecutivo.",
    )


# ── Slide 5: Beneficios del portal de clientes ───────────────────────────────
def s_client_value(prs):
    slide = _blank(prs)
    _title(slide, "Por qué importa el Portal de Clientes")
    _bullets(
        slide,
        [
            "Mejor experiencia: el cliente resuelve solo y a cualquier hora",
            "Menos consultas repetitivas al equipo comercial",
            "Imagen moderna que refuerza los +35 años de trayectoria",
            "Cotizaciones más rápidas, ciclo comercial más corto",
        ],
    )
    _notes(
        slide,
        "Conectar con objetivos de negocio: satisfacción del cliente, eficiencia del equipo "
        "y diferenciación frente a la competencia.",
    )


# ── Slide 6: Portales internos y sus perfiles ────────────────────────────────
def s_internal_roles(prs):
    slide = _blank(prs)
    _title(slide, "Portales Internos", "Cinco perfiles, cada uno ve solo lo suyo")
    roles = [
        ("Administrador", "Gobierno: usuarios, reportes y parámetros"),
        ("Ejecutivo", "Su cartera: clientes, cotizaciones y tracking"),
        ("Pricing", "Dueño de las tarifas y los proveedores"),
        ("Operaciones", "Vista global de todos los clientes"),
        ("Proveedor", "Externo: publica y actualiza sus tarifas"),
    ]
    for i, (t, d) in enumerate(roles):
        row, col = divmod(i, 2)
        _card(slide, 0.65 + col * 6.2, 1.75 + row * 1.5, 5.8, 1.3, t, d, title_size=16, desc_size=12)
    _notes(
        slide,
        "No hace falta explicar cada permiso: la clave es que es un solo sistema que se adapta "
        "al rol de cada persona. Pricing y Ejecutivo pueden combinarse; el resto son exclusivos.",
    )


# ── Slide 7: Qué resuelven los portales internos ─────────────────────────────
def s_internal_value(prs):
    slide = _blank(prs)
    _title(slide, "Qué resuelven los Portales Internos")
    items = [
        ("Cotizar", "El equipo cotiza con tarifas siempre vigentes"),
        ("Clientes", "Cartera y seguimiento comercial ordenados"),
        ("Tarifas", "Proveedores cargan y pricing publica"),
        ("Control", "Reportes, auditoría y gobierno para dirección"),
    ]
    for i, (t, d) in enumerate(items):
        row, col = divmod(i, 2)
        _card(slide, 0.65 + col * 6.2, 1.9 + row * 1.7, 5.8, 1.45, t, d, desc_size=14)
    _notes(
        slide,
        "Mensaje general: reemplaza planillas y correos internos por un flujo ordenado. "
        "Dirección obtiene visibilidad de rentabilidad y desempeño.",
    )


# ── Slide 8: Cómo se conecta todo (flujo) ────────────────────────────────────
def s_flow(prs):
    slide = _blank(prs)
    _title(slide, "Cómo se conecta todo", "Del proveedor al cliente, sin fricción")
    steps = [
        ("Proveedor", "Carga sus tarifas"),
        ("Pricing", "Revisa y publica"),
        ("Equipo/Cliente", "Cotiza al instante"),
        ("Operación", "Se embarca y rastrea"),
    ]
    x = 0.65
    for i, (t, d) in enumerate(steps):
        _card(slide, x, 2.4, 2.7, 1.8, t, d, title_size=17, desc_size=12)
        if i < len(steps) - 1:
            arr = slide.shapes.add_textbox(Inches(x + 2.7), Inches(2.9), Inches(0.55), Inches(0.8))
            ar = arr.text_frame.paragraphs[0].add_run()
            ar.text = "→"
            _set_run(ar, size=28, bold=True, color=PRIMARY)
        x += 3.15
    _notes(
        slide,
        "Esta cadena es el corazón del ecosistema: una tarifa cargada por el proveedor termina "
        "en una cotización que el cliente ve en segundos. Todo trazable.",
    )


# ── Slide 9: Beneficios globales ─────────────────────────────────────────────
def s_benefits(prs):
    slide = _blank(prs)
    _title(slide, "Beneficios para Seemann Group")
    _bullets(
        slide,
        [
            "Información centralizada y consistente en todas las oficinas",
            "Equipos más ágiles: menos trabajo manual, más foco en el cliente",
            "Visibilidad para dirección: rentabilidad, desempeño y cumplimiento",
            "Mejor servicio al cliente y una imagen digital diferenciadora",
        ],
    )
    _notes(
        slide,
        "Cierre de valor antes de los próximos pasos. Repetir la idea de un solo ecosistema "
        "que beneficia a los tres actores a la vez.",
    )


# ── Slide 10: Cierre ─────────────────────────────────────────────────────────
def s_closing(prs):
    slide = _blank(prs)
    _title(slide, "En resumen")
    _bullets(
        slide,
        [
            "Un ecosistema: clientes, equipo interno y proveedores en la misma plataforma",
            "El cliente gana autonomía; el equipo, eficiencia; la dirección, control",
            "Próximos pasos: impulsar adopción y capacitar a cada perfil",
        ],
        top=1.7,
    )
    qa = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), Inches(4.4), Inches(6.0), Inches(1.9)
    )
    qa.fill.solid()
    qa.fill.fore_color.rgb = SECONDARY
    qa.line.fill.background()
    tf = qa.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "¿Preguntas?"
    _set_run(r, size=32, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "portalclientes.seemanngroup.com"
    _set_run(r2, size=14, color=PRIMARY)
    _notes(
        slide,
        "Cerrar con la frase clave: una sola plataforma digital que conecta a todos. "
        "Invitar preguntas.",
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in (
        s_cover, s_ecosystem, s_two_portals, s_client, s_client_value,
        s_internal_roles, s_internal_value, s_flow, s_benefits, s_closing,
    ):
        fn(prs)
    prs.save(str(OUTPUT))
    print(f"Presentación generada: {OUTPUT} ({len(prs.slides)} diapositivas)")


if __name__ == "__main__":
    build()
